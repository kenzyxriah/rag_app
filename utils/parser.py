import os
# pdf
from PyPDF2 import PdfReader
import io
from docx import Document
from pptx import Presentation

class Parser():
    
    def __init__(self):
        pass
    
    # make it async so we can parse multiple files together, concurrently
    async def parse(self, arg: bytes, ext: str):
        modalities = {'pdf': Parser.parse_pdf, 'txt': Parser.parse_txt,
                 'docx': Parser.parse_docx ,'pptx': Parser.parse_pptx}
        
        if ext in modalities:
            content = await modalities[ext](arg)
            return content
        else:
            raise TypeError('Invalid file type')

    
    
    @staticmethod
    async def parse_txt(arg: bytes) -> str:
        try:
            content = arg.decode('utf-8')

            return content
        except Exception as e:
            raise type(e)(f"Error parsing document: {e}") from e
            
    @staticmethod
    async def parse_pdf(file_bytes: bytes, max_allowed_pages=150):
        
        """Parses pdf to texts if not more than 150 pages by default"""
        try:
    
            reader = PdfReader(io.BytesIO(file_bytes))
            total_pages = len(reader.pages)

            if total_pages > max_allowed_pages:
                raise ValueError(f"PDF has {total_pages} pages, which exceeds the limit of {max_allowed_pages}.")

            # If within limits, just extract all text
            text = [
                page.extract_text().strip()
                for page in reader.pages
                if page.extract_text()
            ]
            return '\n'.join(text)

        except Exception as e:
            raise RuntimeError(f"Error parsing PDF: {e}")
        


    @staticmethod
    async def parse_docx(file: bytes):
        try:

            doc = Document(io.BytesIO(file))

            full_text = [para.text for para in doc.paragraphs if para.text]
            
            # Add text from tables as well
            if len(doc.tables) >0:
                for table in doc.tables:
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        if row_text:
                            full_text.append("|\t".join(row_text))
            
            return "\n".join(full_text)

                
        except Exception as e:
            raise type(e)(f"Error parsing document: {e}") from e

    # text only from pptx

    @staticmethod
    async def parse_pptx(file: bytes) -> str:
        try:
            presentation = Presentation(io.BytesIO(file))
            slides_text = []

            for i, slide in enumerate(presentation.slides, start=1):
                slide_text = "\n".join(
                    shape.text.strip()
                    for shape in slide.shapes
                    if shape.has_text_frame and shape.text.strip()
                )
                slides_text.append(f"Page {i}:\n{slide_text}\n")

            return "\n".join(slides_text).strip()
        
        except Exception as e:
            raise
        

